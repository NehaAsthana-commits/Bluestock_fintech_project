"""
Day 6 -- Bluestock Mutual Fund Analytics Dashboard (PowerPoint)
Creates bluestock_mf_dashboard.pptx with 5 slides:
  Slide 0: Title
  Slide 1: Industry Overview   (KPI cards + AUM trend + AUM by AMC)
  Slide 2: Fund Performance    (Risk-Return scatter + scorecard + NAV vs benchmark)
  Slide 3: Investor Analytics  (State bars + donut + age SIP + monthly volume)
  Slide 4: SIP & Market Trends (Dual-axis SIP+NIFTY50 + heatmap + top-5 FY25)
"""

import os, io, warnings
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.ticker as mticker
from matplotlib.dates import DateFormatter as MDateFormatter
from matplotlib.patches import FancyBboxPatch
from matplotlib.backends.backend_pdf import PdfPages
import seaborn as sns
import pandas as pd
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

warnings.filterwarnings("ignore")
plt.rcParams.update({"font.family": "DejaVu Sans"})

BASE = os.path.dirname(os.path.abspath(__file__))
PROC = os.path.join(BASE, "..", "data", "processed")
RAW  = os.path.join(BASE, "..", "data", "raw")
CTMP = os.path.join(BASE, "chart_panels")
os.makedirs(CTMP, exist_ok=True)

# ── Colour constants ─────────────────────────────────────────────────────────
BG_SLIDE  = RGBColor(0xEE, 0xF2, 0xF7)
HDR_COL   = RGBColor(0x0F, 0x17, 0x2A)
PRIMARY   = RGBColor(0x1E, 0x40, 0xAF)
BLUE_R    = RGBColor(0x37, 0x6A, 0xE8)
ORANGE_R  = RGBColor(0xF9, 0x73, 0x16)
GREEN_R   = RGBColor(0x16, 0xA3, 0x4A)
RED_R     = RGBColor(0xDC, 0x26, 0x26)
WHITE_R   = RGBColor(0xFF, 0xFF, 0xFF)
GREY_R    = RGBColor(0x64, 0x74, 0x8B)
LGREY_R   = RGBColor(0xE2, 0xE8, 0xF0)

# matplotlib hex equivalents
BG      = "#EEF2F7"; HEADER = "#0F172A"; PRIM = "#1E40AF"; BLUE = "#2563EB"
ORANGE  = "#F97316"; GREEN  = "#16A34A"; RED  = "#DC2626"; PURPLE = "#7C3AED"
TEAL    = "#0D9488"; GREY   = "#64748B"; LGREY = "#E2E8F0"; GOLD  = "#D97706"
WHITE   = "#FFFFFF"

CAT_PAL = {
    "Large Cap": PRIM, "Mid Cap": BLUE, "Small Cap": ORANGE, "ELSS": GREEN,
    "Flexi Cap": PURPLE, "Large & Mid Cap": TEAL, "Liquid": GREY,
    "Debt": GOLD, "Gilt": RED, "Multi Cap": "#0891B2", "Index": "#059669",
}

# ── Slide dimensions (16:9 widescreen) ──────────────────────────────────────
SW = Inches(13.333)   # slide width
SH = Inches(7.5)      # slide height
DPI_C = 180           # chart PNG resolution

# ── Load data ────────────────────────────────────────────────────────────────
print("Loading data ...")
fm   = pd.read_csv(os.path.join(PROC, "01_fund_master.csv"))
nav  = pd.read_csv(os.path.join(RAW,  "02_nav_history.csv"), parse_dates=["date"])
aum  = pd.read_csv(os.path.join(PROC, "03_aum_by_fund_house.csv"), parse_dates=["date"])
sip  = pd.read_csv(os.path.join(PROC, "04_monthly_sip_inflows.csv"))
cat  = pd.read_csv(os.path.join(PROC, "05_category_inflows.csv"))
fol  = pd.read_csv(os.path.join(PROC, "06_industry_folio_count.csv"))
perf = pd.read_csv(os.path.join(PROC, "07_scheme_performance.csv"))
txn  = pd.read_csv(os.path.join(PROC, "08_investor_transactions.csv"),
                   parse_dates=["transaction_date"])
bm   = pd.read_csv(os.path.join(PROC, "10_benchmark_indices.csv"), parse_dates=["date"])

nav  = nav.sort_values(["amfi_code","date"]).reset_index(drop=True)
sip["month_dt"]  = pd.to_datetime(sip["month"])
cat["month_dt"]  = pd.to_datetime(cat["month"])
txn["month_dt"]  = txn["transaction_date"].dt.to_period("M").dt.to_timestamp()
print("Data loaded.\n")


# ── chart helpers ────────────────────────────────────────────────────────────
def save_chart(fig, name):
    path = os.path.join(CTMP, name)
    fig.savefig(path, dpi=DPI_C, bbox_inches="tight", facecolor=WHITE)
    plt.close(fig)
    return path

def style(ax, title=None, xlabel=None, ylabel=None, grid="y"):
    ax.set_facecolor(WHITE)
    ax.spines[["top","right"]].set_visible(False)
    ax.spines[["left","bottom"]].set_edgecolor(LGREY)
    ax.tick_params(colors=GREY, labelsize=8)
    if grid: ax.grid(axis=grid, color=LGREY, lw=0.6, alpha=0.9, zorder=0)
    if title:  ax.set_title(title, fontsize=10, fontweight="bold", color=HEADER, pad=6)
    if xlabel: ax.set_xlabel(xlabel, fontsize=8.5, color=GREY)
    if ylabel: ax.set_ylabel(ylabel, fontsize=8.5, color=GREY)

def norm100(s):
    return (s / s.dropna().iloc[0]) * 100


# ============================================================
# CHART GENERATION
# ============================================================
print("Generating chart panels ...")

# ── S1: AUM Trend ────────────────────────────────────────────
aum_trend = aum.groupby("date")["aum_crore"].sum().reset_index()
aum_trend["aum_lakh"] = aum_trend["aum_crore"] / 1e5
fig, ax = plt.subplots(figsize=(8.5, 4.4), facecolor=WHITE)
ax.fill_between(aum_trend["date"], aum_trend["aum_lakh"], alpha=0.12, color=PRIM)
ax.plot(aum_trend["date"], aum_trend["aum_lakh"],
        color=PRIM, lw=2.5, marker="o", ms=5, zorder=3)
last = aum_trend.iloc[-1]
ax.annotate(f'₹{last["aum_lakh"]:.1f}L Cr',
            xy=(last["date"], last["aum_lakh"]),
            xytext=(8,6), textcoords="offset points",
            fontsize=9, color=PRIM, fontweight="bold",
            arrowprops=dict(arrowstyle="-", color=LGREY, lw=0.8))
style(ax, "Industry AUM Trend — 2022 to 2026  (10 Fund Houses)",
      ylabel="AUM (₹ Lakh Crore)")
ax.xaxis.set_major_formatter(MDateFormatter("%b '%y"))
ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"₹{v:.0f}L"))
plt.setp(ax.get_xticklabels(), rotation=30, ha="right")
plt.tight_layout()
C_AUM_TREND = save_chart(fig, "s1_aum_trend.png")
print("  s1_aum_trend.png")

# ── S1: AUM by AMC ───────────────────────────────────────────
latest_dt = aum["date"].max()
aum_amc   = (aum[aum["date"]==latest_dt]
             .sort_values("aum_crore")
             .assign(aum_lakh=lambda x: x["aum_crore"]/1e5))
short_names = aum_amc["fund_house"].str.replace(r"\s*Mutual Fund","",regex=True).str[:16]
bar_clr = [PRIM if v==aum_amc["aum_crore"].max() else BLUE for v in aum_amc["aum_crore"]]
fig, ax = plt.subplots(figsize=(4.2, 4.4), facecolor=WHITE)
bars = ax.barh(short_names, aum_amc["aum_lakh"], color=bar_clr, height=0.65)
style(ax, "AUM by AMC — Mar 2026", xlabel="₹ Lakh Crore", grid="x")
ax.spines["left"].set_visible(False)
ax.tick_params(axis="y", labelsize=7.5)
for b in bars:
    ax.text(b.get_width()+0.05, b.get_y()+b.get_height()/2,
            f"₹{b.get_width():.1f}L", va="center", ha="left", fontsize=7.5, color=GREY)
plt.tight_layout()
C_AUM_AMC = save_chart(fig, "s1_aum_amc.png")
print("  s1_aum_amc.png")

# ── S2: Risk-Return Scatter ───────────────────────────────────
fig, ax = plt.subplots(figsize=(5.8, 6.0), facecolor=WHITE)
for cat_name in perf["category"].unique():
    sub  = perf[perf["category"]==cat_name]
    szs  = (sub["aum_crore"]/sub["aum_crore"].max()*380+60).clip(60,500)
    ax.scatter(sub["return_3yr_pct"], sub["std_dev_ann_pct"],
               s=szs, color=CAT_PAL.get(cat_name,GREY),
               alpha=0.78, label=cat_name, edgecolors=WHITE, lw=0.5, zorder=3)
for _,row in perf.nlargest(5,"aum_crore").iterrows():
    nm = row["scheme_name"].split("-")[0].strip()[:16]
    ax.annotate(nm, xy=(row["return_3yr_pct"],row["std_dev_ann_pct"]),
                xytext=(4,3), textcoords="offset points", fontsize=6.5, color=HEADER)
ax.axvline(perf["return_3yr_pct"].mean(), color=LGREY, lw=1, ls="--", alpha=0.8)
ax.axhline(perf["std_dev_ann_pct"].mean(), color=LGREY, lw=1, ls="--", alpha=0.8)
style(ax, "Risk vs Return  (bubble = AUM)",
      xlabel="3-Year Annualised Return (%)", ylabel="Annualised Std Dev (%)", grid=None)
ax.legend(loc="upper left", fontsize=7, framealpha=0.7, markerscale=0.7)
ax.grid(True, color=LGREY, lw=0.4, alpha=0.6)
plt.tight_layout()
C_SCATTER = save_chart(fig, "s2_scatter.png")
print("  s2_scatter.png")

# ── S2: Scorecard Table ───────────────────────────────────────
top10 = perf.nlargest(10,"sharpe_ratio")[
    ["scheme_name","return_3yr_pct","sharpe_ratio","expense_ratio_pct","risk_grade"]
].reset_index(drop=True)
fig, ax = plt.subplots(figsize=(6.5, 3.8), facecolor=WHITE)
ax.set_facecolor(WHITE); ax.set_xlim(0,1); ax.set_ylim(0,1); ax.axis("off")
ax.set_title("Fund Scorecard — Top 10 by Sharpe Ratio",
             fontsize=10, fontweight="bold", color=HEADER, pad=5)
hdrs = ["Scheme Name","3yr Ret%","Sharpe","Exp%","Risk"]
cx   = [0.02, 0.54, 0.67, 0.79, 0.90]
rh   = 0.073
hy   = 0.90
ax.add_patch(plt.Rectangle((0,hy-0.015),1,0.085,fc=PRIM,ec="none",zorder=1))
for j,h in enumerate(hdrs):
    ax.text(cx[j],hy+0.02,h,fontsize=8,color=WHITE,fontweight="bold",va="center")
risk_c = {"Low":GREEN,"Moderate":GOLD,"High":RED,"Very High":RED,"Moderately High":ORANGE}
for i,row in top10.iterrows():
    y = hy-(i+1)*rh
    ax.add_patch(plt.Rectangle((0,y-0.01),1,rh,fc=LGREY if i%2==0 else WHITE,ec="none",zorder=0))
    nm = row["scheme_name"].split("-")[0].strip()[:32]
    ax.text(cx[0],y+0.022,nm,fontsize=7,color=HEADER,va="center")
    ax.text(cx[1],y+0.022,f"{row['return_3yr_pct']:.1f}%",fontsize=7.5,color=PRIM,fontweight="bold",va="center")
    ax.text(cx[2],y+0.022,f"{row['sharpe_ratio']:.2f}",fontsize=7.5,color=HEADER,va="center")
    ax.text(cx[3],y+0.022,f"{row['expense_ratio_pct']:.2f}%",fontsize=7.5,color=HEADER,va="center")
    rc = risk_c.get(str(row["risk_grade"]),GREY)
    ax.text(cx[4],y+0.022,str(row["risk_grade"])[:3],fontsize=7.5,color=rc,fontweight="bold",va="center")
plt.tight_layout()
C_SCORECARD = save_chart(fig, "s2_scorecard.png")
print("  s2_scorecard.png")

# ── S2: NAV vs NIFTY100 ───────────────────────────────────────
top3 = (perf[perf["category"]=="Equity"]
        .sort_values("aum_crore",ascending=False)
        .head(3)[["amfi_code","scheme_name"]])
nifty100 = bm[bm["index_name"]=="NIFTY100"].set_index("date")["close_value"]
fig, ax = plt.subplots(figsize=(6.5, 3.1), facecolor=WHITE)
nav_cols = [GREEN, TEAL, PURPLE]
for (_, fr), col in zip(top3.iterrows(), nav_cols):
    fn = nav[nav["amfi_code"]==fr["amfi_code"]].set_index("date")["nav"]
    ax.plot(norm100(fn).index, norm100(fn).values, lw=1.8, color=col, alpha=0.9,
            label=fr["scheme_name"].split("-")[0].strip()[:22])
n_idx = pd.date_range(nifty100.index.min(), nifty100.index.max(), freq="B")
nifty_n = norm100(nifty100.reindex(n_idx).ffill())
ax.plot(nifty_n.index, nifty_n.values, lw=2.0, color=ORANGE, ls="--", label="NIFTY 100")
style(ax, "NAV vs NIFTY 100  (Indexed to 100 at start)", ylabel="Indexed Value")
ax.xaxis.set_major_formatter(MDateFormatter("%b '%y"))
plt.setp(ax.get_xticklabels(), rotation=25, ha="right")
ax.legend(loc="upper left", fontsize=7, framealpha=0.7, ncol=2)
plt.tight_layout()
C_NAV = save_chart(fig, "s2_nav.png")
print("  s2_nav.png")

# ── S3: State Transaction Bars ────────────────────────────────
state_cr = (txn.groupby("state")["amount_inr"].sum()
            .sort_values().tail(15) / 1e7)
fig, ax = plt.subplots(figsize=(7.0, 3.8), facecolor=WHITE)
clrs = [PRIM if i==len(state_cr)-1 else BLUE for i in range(len(state_cr))]
ax.barh(state_cr.index, state_cr.values, color=clrs, height=0.68)
style(ax, "Transaction Amount by State — Top 15", xlabel="₹ Crore", grid="x")
ax.spines["left"].set_visible(False)
ax.tick_params(axis="y", labelsize=7.5)
for i,v in enumerate(state_cr.values):
    ax.text(v+state_cr.max()*0.01, i, f"₹{v:.0f}Cr",
            va="center", ha="left", fontsize=7, color=GREY)
plt.tight_layout()
C_STATE = save_chart(fig, "s3_state.png")
print("  s3_state.png")

# ── S3: Donut ─────────────────────────────────────────────────
donut_data = txn.groupby("transaction_type")["amount_inr"].sum()
dcols = {
    "SIP":BLUE,"Lumpsum":ORANGE,"Redemption":RED,
    "STP":PURPLE,"SWP":TEAL,"Switch":GREY,
}
fig, ax = plt.subplots(figsize=(5.0, 3.8), facecolor=WHITE)
wedges,_,autotexts = ax.pie(
    donut_data.values, labels=None,
    colors=[dcols.get(k,GREY) for k in donut_data.index],
    autopct="%1.1f%%", startangle=90, pctdistance=0.78,
    wedgeprops=dict(width=0.52, edgecolor=WHITE, linewidth=2))
for at in autotexts:
    at.set_fontsize(8.5); at.set_color(WHITE); at.set_fontweight("bold")
legend_p = [mpatches.Patch(color=dcols.get(k,GREY),
            label=f"{k}  ₹{v/1e7:.0f}Cr") for k,v in donut_data.items()]
ax.legend(handles=legend_p, loc="lower center", bbox_to_anchor=(0.5,-0.15),
          fontsize=7.5, ncol=2, framealpha=0.5)
ax.set_title("SIP / Lumpsum / Redemption Split", fontsize=10,
             fontweight="bold", color=HEADER, pad=5)
plt.tight_layout()
C_DONUT = save_chart(fig, "s3_donut.png")
print("  s3_donut.png")

# ── S3: Age Group vs Avg SIP ──────────────────────────────────
age_sip = (txn[txn["transaction_type"]=="SIP"]
           .groupby("age_group")["amount_inr"].mean())
age_order = [a for a in ["18-25","26-35","36-45","46-55","56+"] if a in age_sip.index]
age_data  = age_sip.reindex(age_order)
fig, ax = plt.subplots(figsize=(5.5, 3.1), facecolor=WHITE)
bars_a = ax.bar(age_data.index, age_data.values,
                color=[BLUE,GREEN,ORANGE,PURPLE,TEAL][:len(age_data)], width=0.58)
style(ax, "Age Group vs Avg SIP Amount", xlabel="Age Group", ylabel="Avg SIP (₹)")
ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"₹{v:,.0f}"))
for b in bars_a:
    ax.text(b.get_x()+b.get_width()/2, b.get_height()+age_data.max()*0.02,
            f"₹{b.get_height():,.0f}", ha="center", va="bottom",
            fontsize=8, color=HEADER, fontweight="bold")
plt.tight_layout()
C_AGE = save_chart(fig, "s3_age.png")
print("  s3_age.png")

# ── S3: Monthly Transaction Volume ────────────────────────────
monthly_vol = (txn.groupby("month_dt")["amount_inr"].sum() / 1e7).reset_index()
fig, ax = plt.subplots(figsize=(6.5, 3.1), facecolor=WHITE)
ax.fill_between(monthly_vol["month_dt"], monthly_vol["amount_inr"],
                alpha=0.12, color=TEAL)
ax.plot(monthly_vol["month_dt"], monthly_vol["amount_inr"],
        color=TEAL, lw=2.0, marker="o", ms=3.5, zorder=3)
style(ax, "Monthly Transaction Volume", ylabel="₹ Crore")
ax.xaxis.set_major_formatter(MDateFormatter("%b '%y"))
ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"₹{v:.0f}Cr"))
plt.setp(ax.get_xticklabels(), rotation=30, ha="right")
plt.tight_layout()
C_MVOL = save_chart(fig, "s3_mvol.png")
print("  s3_mvol.png")

# ── S4: Dual-Axis SIP + NIFTY50 ──────────────────────────────
nifty50 = (bm[bm["index_name"]=="NIFTY50"].set_index("date")["close_value"]
           .resample("ME").last().reset_index().rename(columns={"date":"dt","close_value":"nifty"}))
sip_m   = sip[["month_dt","sip_inflow_crore"]].sort_values("month_dt")
merged  = pd.merge_asof(sip_m, nifty50.rename(columns={"dt":"month_dt"}),
                        on="month_dt", direction="nearest")
fig, ax = plt.subplots(figsize=(12.5, 3.8), facecolor=WHITE)
ax2 = ax.twinx()
ax.bar(merged["month_dt"], merged["sip_inflow_crore"],
       width=20, color=BLUE, alpha=0.75, label="SIP Inflow (₹ Cr)", zorder=2)
ax2.plot(merged["month_dt"], merged["nifty"], color=ORANGE, lw=2.2,
         marker="o", ms=3.5, label="NIFTY 50", zorder=3)
ax.set_ylabel("SIP Inflow (₹ Crore)", fontsize=9, color=BLUE)
ax2.set_ylabel("NIFTY 50 Level", fontsize=9, color=ORANGE)
ax.tick_params(axis="y", colors=BLUE, labelsize=8)
ax2.tick_params(axis="y", colors=ORANGE, labelsize=8)
ax.set_facecolor(WHITE)
ax.spines[["top"]].set_visible(False)
ax.spines[["left","bottom"]].set_edgecolor(LGREY)
ax2.spines[["top"]].set_visible(False)
ax.grid(axis="y", color=LGREY, lw=0.5, alpha=0.7, zorder=0)
ax.xaxis.set_major_formatter(MDateFormatter("%b '%y"))
plt.setp(ax.get_xticklabels(), rotation=30, ha="right", fontsize=7.5)
ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"₹{v:,.0f}"))
ax2.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"{v:,.0f}"))
ax.set_title("Monthly SIP Inflows vs NIFTY 50  (2022–2026)",
             fontsize=10, fontweight="bold", color=HEADER, pad=5)
handles = [mpatches.Patch(fc=BLUE,label="SIP Inflow (₹ Cr)"),
           mpatches.Patch(fc=ORANGE,label="NIFTY 50")]
ax.legend(handles=handles, loc="upper left", fontsize=8.5, framealpha=0.7)
plt.tight_layout()
C_DUAL = save_chart(fig, "s4_dual.png")
print("  s4_dual.png")

# ── S4: Category Heatmap ─────────────────────────────────────
cat_rec = cat[cat["month_dt"] >= cat["month_dt"].max()-pd.DateOffset(months=11)]
hm = cat_rec.pivot_table(index="category", columns="month_dt",
                          values="net_inflow_crore", aggfunc="sum")
hm.columns = [d.strftime("%b'%y") for d in hm.columns]
fig, ax = plt.subplots(figsize=(7.5, 3.6), facecolor=WHITE)
sns.heatmap(hm, ax=ax, cmap="YlOrRd", linewidths=0.4, linecolor=BG,
            fmt=".0f", annot=True, annot_kws={"size":7},
            cbar_kws={"shrink":0.75,"pad":0.01})
ax.set_title("Category Inflow Heatmap — Last 12 Months (₹ Cr)",
             fontsize=10, fontweight="bold", color=HEADER, pad=5)
ax.set_xlabel(""); ax.set_ylabel("")
ax.tick_params(axis="x", rotation=40, labelsize=7)
ax.tick_params(axis="y", rotation=0,  labelsize=7.5)
plt.tight_layout()
C_HEATMAP = save_chart(fig, "s4_heatmap.png")
print("  s4_heatmap.png")

# ── S4: Top 5 Categories FY25 ────────────────────────────────
fy25     = cat[(cat["month_dt"]>="2024-04-01") & (cat["month_dt"]<="2025-03-31")]
top5_cat = fy25.groupby("category")["net_inflow_crore"].sum().nlargest(5).sort_values()
fig, ax  = plt.subplots(figsize=(5.0, 3.6), facecolor=WHITE)
clrs_t5  = [CAT_PAL.get(c,BLUE) for c in top5_cat.index]
bars_t5  = ax.barh(top5_cat.index, top5_cat.values/1e3, color=clrs_t5, height=0.58)
style(ax, "Top 5 Categories — Net Inflow FY25", xlabel="₹ '000 Crore", grid="x")
ax.spines["left"].set_visible(False)
ax.tick_params(axis="y", labelsize=8.5)
for b in bars_t5:
    ax.text(b.get_width()+top5_cat.max()/1e3*0.02,
            b.get_y()+b.get_height()/2,
            f"₹{b.get_width():.1f}K Cr",
            va="center", ha="left", fontsize=8.5, color=HEADER, fontweight="bold")
plt.tight_layout()
C_TOP5 = save_chart(fig, "s4_top5.png")
print("  s4_top5.png\n")


# ============================================================
# PPT ASSEMBLY
# ============================================================
print("Building PowerPoint ...")

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
blank = prs.slide_layouts[6]   # truly blank layout

# ── helpers ──────────────────────────────────────────────────
def set_bg(slide, color_rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=12, bold=False,
             color=WHITE_R, align=PP_ALIGN.LEFT, wrap=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.text_frame.word_wrap = wrap
    p = txb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_header_bar(slide, page_title):
    bar = add_rect(slide, 0, 0, 13.333, 0.58, HDR_COL)
    add_text(slide, "BLUESTOCK", 0.18, 0.05, 1.8, 0.48, size=18,
             bold=True, color=RGBColor(0x93,0xC5,0xFD))
    add_text(slide, "|", 1.88, 0.05, 0.3, 0.48, size=18, color=GREY_R)
    add_text(slide, page_title, 2.12, 0.06, 7.0, 0.48, size=17,
             bold=True, color=WHITE_R)
    add_text(slide, "Mutual Fund Analytics  |  Bluestock Fintech Capstone",
             7.8, 0.16, 5.3, 0.3, size=8, color=GREY_R, align=PP_ALIGN.RIGHT)

def add_pic(slide, path, x, y, w, h):
    slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))

def add_kpi_card(slide, x, y, w, h, value, label, sub, accent_rgb):
    # White card background
    card = add_rect(slide, x, y, w, h, WHITE_R, line_rgb=accent_rgb)
    # Accent top strip
    add_rect(slide, x+0.02, y+0.02, w-0.04, 0.26, accent_rgb)
    # Label (in accent strip)
    add_text(slide, label, x+0.04, y+0.03, w-0.08, 0.22,
             size=8, bold=True, color=WHITE_R, align=PP_ALIGN.CENTER)
    # Value
    add_text(slide, value, x+0.04, y+0.31, w-0.08, 0.52,
             size=20, bold=True, color=RGBColor(0x0F,0x17,0x2A),
             align=PP_ALIGN.CENTER)
    # Sub-label
    add_text(slide, sub, x+0.04, y+0.85, w-0.08, 0.20,
             size=7.5, bold=False, color=GREY_R, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────
# SLIDE 0 — Title slide
# ─────────────────────────────────────────────────────────────
slide0 = prs.slides.add_slide(blank)
set_bg(slide0, HDR_COL)
add_rect(slide0, 0, 0, 13.333, 7.5, HDR_COL)

# Blue accent band
add_rect(slide0, 0, 2.6, 13.333, 0.06, RGBColor(0x1E,0x40,0xAF))
add_rect(slide0, 0, 5.1, 13.333, 0.06, RGBColor(0x1E,0x40,0xAF))

add_text(slide0, "BLUESTOCK FINTECH", 1.5, 0.8, 10.0, 1.0,
         size=36, bold=True, color=RGBColor(0x93,0xC5,0xFD), align=PP_ALIGN.CENTER)
add_text(slide0, "MUTUAL FUND ANALYTICS DASHBOARD", 1.0, 1.75, 11.0, 0.9,
         size=26, bold=True, color=WHITE_R, align=PP_ALIGN.CENTER)
add_text(slide0, "Capstone Project I  |  Group: asthananeha2015  /  bsoham2704  /  Bubai Das",
         1.5, 2.8, 10.0, 0.5, size=12, color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.CENTER)
add_text(slide0, "Day 6 — Dashboard Development", 2.5, 3.4, 8.0, 0.5,
         size=12, color=RGBColor(0x64,0x74,0x8B), align=PP_ALIGN.CENTER)
add_text(slide0,
         "Page 1: Industry Overview   |   Page 2: Fund Performance\n"
         "Page 3: Investor Analytics   |   Page 4: SIP & Market Trends",
         2.0, 4.1, 9.0, 1.0, size=11, color=RGBColor(0x94,0xA3,0xB8),
         align=PP_ALIGN.CENTER, wrap=True)
add_text(slide0, "Powered by Python  |  Data: AMFI / MFAPI / NSE / BSE",
         2.0, 5.3, 9.0, 0.4, size=9.5, color=RGBColor(0x47,0x55,0x69),
         align=PP_ALIGN.CENTER)
add_text(slide0, "Bluestock Fintech  |  July 2026",
         3.0, 6.5, 7.0, 0.45, size=10, color=RGBColor(0x47,0x55,0x69),
         align=PP_ALIGN.CENTER)
print("  Slide 0: Title")


# ─────────────────────────────────────────────────────────────
# SLIDE 1 — Industry Overview
# ─────────────────────────────────────────────────────────────
slide1 = prs.slides.add_slide(blank)
set_bg(slide1, RGBColor(0xEE,0xF2,0xF7))
add_header_bar(slide1, "Industry Overview")

kpi_specs = [
    ("₹81L Cr",  "TOTAL INDUSTRY AUM",  "Mar 2026",  RGBColor(0x1E,0x40,0xAF)),
    ("₹31K Cr",  "ANNUAL SIP INFLOWS",  "FY2025",    RGBColor(0x37,0x6A,0xE8)),
    ("26.12 Cr", "TOTAL FOLIOS",        "Q4 FY2026", RGBColor(0x16,0xA3,0x4A)),
    ("1,908",    "TOTAL SCHEMES",       "Active",    RGBColor(0xF9,0x73,0x16)),
]
card_w, card_h = 3.02, 1.25
gap = 0.087
start_x = (13.333 - 4*card_w - 3*gap) / 2
for i,(val,lbl,sub,col) in enumerate(kpi_specs):
    cx = start_x + i*(card_w+gap)
    add_kpi_card(slide1, cx, 0.68, card_w, card_h, val, lbl, sub, col)

add_pic(slide1, C_AUM_TREND, 0.15, 2.08, 8.85, 5.20)
add_pic(slide1, C_AUM_AMC,   9.10, 2.08, 4.05, 5.20)
print("  Slide 1: Industry Overview")


# ─────────────────────────────────────────────────────────────
# SLIDE 2 — Fund Performance
# ─────────────────────────────────────────────────────────────
slide2 = prs.slides.add_slide(blank)
set_bg(slide2, RGBColor(0xEE,0xF2,0xF7))
add_header_bar(slide2, "Fund Performance")

add_text(slide2, "Slicers: Fund House  |  Category  |  Plan",
         8.5, 0.62, 4.6, 0.3, size=7.5, color=GREY_R, align=PP_ALIGN.RIGHT)

add_pic(slide2, C_SCATTER,  0.15, 0.65, 6.05, 6.65)
add_pic(slide2, C_SCORECARD, 6.4, 0.65, 6.75, 3.80)
add_pic(slide2, C_NAV,       6.4, 4.55, 6.75, 2.75)
print("  Slide 2: Fund Performance")


# ─────────────────────────────────────────────────────────────
# SLIDE 3 — Investor Analytics
# ─────────────────────────────────────────────────────────────
slide3 = prs.slides.add_slide(blank)
set_bg(slide3, RGBColor(0xEE,0xF2,0xF7))
add_header_bar(slide3, "Investor Analytics")

add_text(slide3, "Slicers: State  |  Age Group  |  City Tier",
         9.5, 0.62, 3.6, 0.3, size=7.5, color=GREY_R, align=PP_ALIGN.RIGHT)

add_pic(slide3, C_STATE, 0.15, 0.65, 7.20, 3.60)
add_pic(slide3, C_DONUT, 7.50, 0.65, 5.65, 3.60)
add_pic(slide3, C_AGE,   0.15, 4.35, 5.90, 3.00)
add_pic(slide3, C_MVOL,  6.20, 4.35, 6.95, 3.00)
print("  Slide 3: Investor Analytics")


# ─────────────────────────────────────────────────────────────
# SLIDE 4 — SIP & Market Trends
# ─────────────────────────────────────────────────────────────
slide4 = prs.slides.add_slide(blank)
set_bg(slide4, RGBColor(0xEE,0xF2,0xF7))
add_header_bar(slide4, "SIP & Market Trends")

add_pic(slide4, C_DUAL,    0.15, 0.65, 12.95, 3.50)
add_pic(slide4, C_HEATMAP, 0.15, 4.25,  7.70, 3.10)
add_pic(slide4, C_TOP5,    8.10, 4.25,  5.05, 3.10)
print("  Slide 4: SIP & Market Trends")


# ─────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────
out = os.path.join(BASE, "bluestock_mf_dashboard.pptx")
prs.save(out)
print(f"\nSaved: {out}")
print(f"Slides: {len(prs.slides)}  (Title + 4 dashboard pages)")
print("\nDay 6 Dashboard (PPT) Complete!")
